"""
会话导出 PPT —— 把对话转成演示文稿

策略:
- 读取 sessions/auto/<file>.json 的 messages
- 提取每轮对话(user + assistant)作为一张幻灯片
- 标题 = 会话日期 / 第一条 user 消息
- 用户消息 → "需求"幻灯片
- AI 回答 → "回答"幻灯片(支持 Markdown 简化)
- 关键工具调用 → "步骤"幻灯片

依赖:python-pptx(可选,无则回退到 Markdown 大纲)
"""
import os
import re
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple

from fr_cli.core.store import JsonStore


DEFAULT_TEMPLATE_COLORS = {
    "title_bg": "1F4E79",
    "title_fg": "FFFFFF",
    "user_bg": "F2F2F2",
    "user_fg": "333333",
    "ai_bg": "DEEBF7",
    "ai_fg": "1F4E79",
}


def load_session_messages(session_path: str) -> List[Dict[str, Any]]:
    """加载会话消息"""
    if not os.path.exists(session_path):
        return []
    try:
        # 兼容 v2 增量格式
        try:
            from fr_cli.memory.incremental import read_session_full_data
            return read_session_full_data(session_path).get("messages", [])
        except Exception:
            pass
        data = JsonStore(session_path, default=dict).read()
        return data.get("messages", [])
    except Exception:
        return []


def extract_conversation_pairs(messages: List[Dict[str, Any]]) -> List[Tuple[Dict[str, Any], Dict[str, Any]]]:
    """提取 user + assistant 配对

    Returns:
        [(user_msg, assistant_msg), ...]
    """
    pairs = []
    i = 0
    while i < len(messages):
        if messages[i].get("role") == "user":
            user = messages[i]
            assistant = None
            if i + 1 < len(messages) and messages[i + 1].get("role") == "assistant":
                assistant = messages[i + 1]
                i += 2
            else:
                i += 1
            pairs.append((user, assistant))
        else:
            i += 1
    return pairs


def clean_markdown_for_ppt(text: str, max_length: int = 1500) -> str:
    """清理 markdown 用于 PPT 文本

    - 移除代码块(保留语言提示)
    - 限制长度
    - 简化链接
    """
    if not text:
        return ""

    # 移除代码块 ```...```
    text = re.sub(r"```[\s\S]*?```", "[代码块已省略]", text)
    # 行内代码
    text = re.sub(r"`([^`]+)`", r"\1", text)
    # 简化链接 [text](url) → text
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    # 移除标题符号
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)
    # 移除粗体/斜体标记
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", text)

    if len(text) > max_length:
        text = text[:max_length] + "..."
    return text.strip()


def build_outline(messages: List[Dict[str, Any]], title: Optional[str] = None,
                  max_slides: int = 50) -> List[Dict[str, Any]]:
    """构建 PPT 大纲

    Returns:
        [
            {"type": "cover", "title": str, "subtitle": str},
            {"type": "user", "content": str, "turn": int},
            {"type": "ai", "content": str, "turn": int},
            ...
            {"type": "summary", "content": str}
        ]
    """
    pairs = extract_conversation_pairs(messages)

    outline = []

    # 封面
    first_user = pairs[0][0].get("content", "未命名会话") if pairs else "空会话"
    title_text = title or (first_user[:50] + ("..." if len(first_user) > 50 else ""))
    subtitle_text = datetime.now().strftime("%Y-%m-%d %H:%M")

    outline.append({
        "type": "cover",
        "title": title_text,
        "subtitle": subtitle_text,
    })

    # 内容:每对 = 2 张幻灯片(user + ai)
    for i, (user, assistant) in enumerate(pairs[:max_slides]):
        # 用户
        user_content = clean_markdown_for_ppt(user.get("content", ""))
        if user_content:
            outline.append({
                "type": "user",
                "content": user_content,
                "turn": i + 1,
                "title": f"需求 #{i+1}",
            })
        # AI
        if assistant:
            ai_content = clean_markdown_for_ppt(assistant.get("content", ""))
            if ai_content:
                outline.append({
                    "type": "ai",
                    "content": ai_content,
                    "turn": i + 1,
                    "title": f"回答 #{i+1}",
                })

    # 总结
    outline.append({
        "type": "summary",
        "content": f"共 {len(pairs)} 轮对话,{len(messages)} 条消息",
        "title": "总结",
    })

    return outline


def export_to_pptx(outline: List[Dict[str, Any]], output_path: str,
                   colors: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """导出为 PPTX

    Args:
        outline: PPT 大纲
        output_path: 输出 .pptx 路径
        colors: 颜色配置

    Returns:
        {"ok": bool, "path": str, "slides": int, "error": str?}
    """
    colors = colors or DEFAULT_TEMPLATE_COLORS

    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
    except ImportError:
        return {"ok": False, "error": "需要 python-pptx:pip install python-pptx"}

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for item in outline:
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            if item["type"] == "cover":
                # 封面:大标题 + 副标题
                _add_text_box(
                    slide, item["title"],
                    left=Inches(1), top=Inches(2.5), width=Inches(11.33), height=Inches(2),
                    font_size=44, bold=True,
                    color=RGBColor.from_string(colors["title_fg"]),
                    bg_color=RGBColor.from_string(colors["title_bg"]),
                )
                _add_text_box(
                    slide, item["subtitle"],
                    left=Inches(1), top=Inches(5), width=Inches(11.33), height=Inches(1),
                    font_size=24,
                    color=RGBColor.from_string(colors["title_fg"]),
                    bg_color=RGBColor.from_string(colors["title_bg"]),
                )
            elif item["type"] in ("user", "ai"):
                bg = colors["user_bg"] if item["type"] == "user" else colors["ai_bg"]
                fg = colors["user_fg"] if item["type"] == "user" else colors["ai_fg"]
                role_label = "👤 用户" if item["type"] == "user" else "🤖 AI"
                # 标题栏
                _add_text_box(
                    slide, f"{role_label} - {item.get('title', '')}",
                    left=Inches(0.5), top=Inches(0.3), width=Inches(12.33), height=Inches(0.7),
                    font_size=20, bold=True,
                    color=RGBColor.from_string(fg),
                    bg_color=RGBColor.from_string(bg),
                )
                # 内容
                _add_text_box(
                    slide, item["content"],
                    left=Inches(0.5), top=Inches(1.2), width=Inches(12.33), height=Inches(5.8),
                    font_size=18,
                    color=RGBColor.from_string(fg),
                    bg_color=RGBColor.from_string("FFFFFF"),
                )
            elif item["type"] == "summary":
                _add_text_box(
                    slide, item["title"],
                    left=Inches(1), top=Inches(2), width=Inches(11.33), height=Inches(1.5),
                    font_size=44, bold=True,
                    color=RGBColor.from_string(colors["title_fg"]),
                    bg_color=RGBColor.from_string(colors["title_bg"]),
                )
                _add_text_box(
                    slide, item["content"],
                    left=Inches(2), top=Inches(4.5), width=Inches(9.33), height=Inches(1),
                    font_size=24,
                    color=RGBColor.from_string(colors["ai_fg"]),
                    bg_color=RGBColor.from_string(colors["ai_bg"]),
                )

        prs.save(output_path)
        return {"ok": True, "path": output_path, "slides": len(outline)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def _add_text_box(slide, text, left, top, width, height,
                  font_size=18, bold=False, color=None, bg_color=None):
    """加文本框"""
    from pptx.util import Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

    if bg_color:
        fill = box.fill
        fill.solid()
        fill.fore_color.rgb = bg_color


def export_to_markdown(outline: List[Dict[str, Any]], output_path: str) -> Dict[str, Any]:
    """降级:导出为 Markdown 大纲(无 python-pptx 时)"""
    try:
        lines = []
        for item in outline:
            if item["type"] == "cover":
                lines.append(f"# {item['title']}")
                lines.append(f"\n*{item['subtitle']}*\n")
            elif item["type"] == "user":
                lines.append(f"## 👤 {item.get('title', '用户')}\n")
                lines.append(item["content"])
                lines.append("")
            elif item["type"] == "ai":
                lines.append(f"## 🤖 {item.get('title', 'AI')}\n")
                lines.append(item["content"])
                lines.append("")
            elif item["type"] == "summary":
                lines.append(f"---\n# {item['title']}\n\n{item['content']}\n")

        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        return {"ok": True, "path": output_path, "format": "markdown", "slides": len(outline)}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def export_session_to_ppt(session_path: str, output_path: Optional[str] = None,
                           format: str = "auto",
                           title: Optional[str] = None,
                           max_slides: int = 50) -> Dict[str, Any]:
    """统一导出入口

    Args:
        session_path: 会话 JSON 路径
        output_path: 输出路径(默认 ~/.fr_cli/exports/session_<timestamp>.pptx)
        format: "pptx" / "markdown" / "auto"(优先 pptx,失败 fallback)
        title: PPT 标题(默认用第一条 user 消息)
        max_slides: 最大幻灯片数

    Returns:
        {"ok": bool, "path": str, "format": str, "slides": int, "error": str?}
    """
    messages = load_session_messages(session_path)
    if not messages:
        return {"ok": False, "error": f"无法读取会话或会话为空: {session_path}"}

    outline = build_outline(messages, title=title, max_slides=max_slides)

    # 默认输出路径
    if not output_path:
        from fr_cli.conf.paths import ROOT as FR_CLI_DIR
        export_dir = FR_CLI_DIR / "exports"
        export_dir.mkdir(parents=True, exist_ok=True)
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        ext = "pptx" if format == "pptx" else "md"
        output_path = str(export_dir / f"session_{ts}.{ext}")

    # 尝试 PPTX
    if format in ("pptx", "auto"):
        result = export_to_pptx(outline, output_path)
        if result["ok"]:
            return {**result, "format": "pptx"}
        if format == "pptx":
            return result
        # auto:fallback to markdown
        output_md = output_path.rsplit(".", 1)[0] + ".md"
        result_md = export_to_markdown(outline, output_md)
        if result_md["ok"]:
            return {**result_md, "pptx_error": result.get("error")}
        return result

    # markdown
    return export_to_markdown(outline, output_path)
